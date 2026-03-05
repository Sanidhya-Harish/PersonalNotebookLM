import os
from pptx import Presentation
from pptx.util import Inches

def create_slides(notebook_id, content_list):

    # Ensure storage folder exists
    os.makedirs("storage", exist_ok=True)

    prs = Presentation()

    # Title slide (use notebook name as title)
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = f"{notebook_id} Presentation"

    # Content slides
    for slide_data in content_list:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_data["title"]

        body = slide.placeholders[1]
        body.text = "\n".join(slide_data["points"])

    # One file per notebook (overwrite)
    filename = f"{notebook_id}_slides.pptx"
    file_path = os.path.join("storage", filename)

    prs.save(file_path)

    return f"storage/{filename}"